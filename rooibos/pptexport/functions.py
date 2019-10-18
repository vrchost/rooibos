from PIL import Image
from django.conf import settings
from pptx import Presentation
from pptx.dml.color import RGBColor

from rooibos.storage.functions import get_image_for_record


SLIDE_LAYOUT_TITLE = 0
SLIDE_LAYOUT_TITLE_AND_CONTENT = 1
SLIDE_LAYOUT_TWO_COLUMNS = 3
SLIDE_LAYOUT_BLANK = 6


COLORS = {
    'white': RGBColor(255, 255, 255),
    'grey': RGBColor(128, 128, 128),
    'black': RGBColor(0, 0, 0),
}

FONT_COLORS = {
    'white': RGBColor(0, 0, 0),
    'grey': RGBColor(0, 0, 0),
    'black': RGBColor(255, 255, 255),
}


class PowerPointGenerator:

    def __init__(self, presentation, user):
        self.presentation = presentation
        self.items = presentation.items.filter(hidden=False)
        self.user = user

    def generate(self, outfile, color, titles, metadata):
        if len(self.items) == 0:
            return False
        pptx_file = Presentation()
        fill = pptx_file.slide_master.background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS[color]
        self.text_color = FONT_COLORS[color]
        self.make_title_slide(pptx_file)
        self.process_slides(pptx_file, titles, metadata)
        pptx_file.save(outfile)
        return True

    def make_title_slide(self, pptx_file):
        slide_layout = pptx_file.slide_layouts[SLIDE_LAYOUT_TITLE]
        slide = pptx_file.slides.add_slide(slide_layout)
        tf = slide.shapes.title.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb = self.text_color
        p.text = self.presentation.title

    def process_slides(self, pptx_file, titles, metadata):
        for item in self.items:
            try:
                image = get_image_for_record(
                    item.record,
                    self.user,
                    getattr(settings, 'PPTEXPORT_WIDTH', 800),
                    getattr(settings, 'PPTEXPORT_HEIGHT', 600)
                )
            except:
                image = None
            values = item.get_fieldvalues()
            if metadata:
                layout = SLIDE_LAYOUT_TWO_COLUMNS
            elif titles:
                layout = SLIDE_LAYOUT_TITLE_AND_CONTENT
            else:
                layout = SLIDE_LAYOUT_BLANK
            slide_layout = pptx_file.slide_layouts[layout]
            slide = pptx_file.slides.add_slide(slide_layout)
            self.add_metadata_to_textframe(
                slide.notes_slide.notes_text_frame, values)
            if titles:
                tf = slide.shapes.title.text_frame
                p = tf.paragraphs[0]
                p.font.color.rgb = self.text_color
                p.text = item.title
            if not image:
                continue

            if not metadata and not titles:
                self.insert_image(
                    slide, image, 0, 0,
                    pptx_file.slide_width, pptx_file.slide_height
                )
                continue

            def insert_image(placeholder):
                self.insert_image(
                    slide, image, placeholder.left,
                    placeholder.top, placeholder.width,
                    placeholder.height
                )

            def insert_metadata(placeholder):
                self.insert_metadata(
                    slide, values, placeholder.left,
                    placeholder.top, placeholder.width,
                    placeholder.height
                )

            methods = [insert_image]
            if metadata:
                methods.append(insert_metadata)

            for shape in slide.shapes:
                if methods and shape.is_placeholder and str(
                        shape.placeholder_format.type).startswith('OBJECT'):
                    methods.pop(0)(
                        slide.placeholders[shape.placeholder_format.idx]
                    )

    def insert_image(self, slide, image, left, top, width, height):
        iwidth, iheight = Image.open(image).size
        ratio = iwidth / iheight
        if ratio > (width / height):
            # need to center vertically
            space = (height - width / ratio) / 2
            slide.shapes.add_picture(
                image, left, top + space, width, height - space * 2)
        else:
            # need to center horizontally
            space = (width - height * ratio) / 2
            slide.shapes.add_picture(
                image, left + space, top, width - space * 2, height)

    def insert_metadata(self, slide, values, left, top, width, height):
        box = slide.shapes.add_textbox(left, top, width, height)
        self.add_metadata_to_textframe(box.text_frame, values)

    def add_metadata_to_textframe(self, text_frame, values):
        p = text_frame.paragraphs[0]
        for value in values:
            run = p.add_run()
            run.font.color.rgb = self.text_color
            run.text = '%s: %s' % (value.resolved_label, value.value)
            p = text_frame.add_paragraph()
        text_frame.word_wrap = True
